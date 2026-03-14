from __future__ import annotations

import base64
import tempfile
from pathlib import Path
from unittest import TestCase
from unittest.mock import patch

import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from office_agent.renderer import PresentationRenderer
from office_agent.schema import PresentationSpec


PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9WnJyoAAAAAASUVORK5CYII="
)


class RendererTests(TestCase):
    def test_renderer_creates_pptx_with_supported_new_slide_types(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "theme": {
                    "preset": "editorial",
                    "custom": {"primary_color": "#224466", "cover_style": "minimal"},
                },
                "slides": [
                    {"type": "title", "part": "opening", "title": "Intro", "subtitle": "Subtitle"},
                    {
                        "type": "timeline",
                        "part": "content",
                        "title": "Roadmap",
                        "events": [
                            {"label": "Q1", "title": "Plan", "detail": "Scope"},
                            {"label": "Q2", "title": "Ship", "detail": "Launch"},
                        ],
                    },
                    {"type": "quote", "part": "content", "title": "Voice", "quote": "Stay close to users.", "attribution": "PM"},
                    {
                        "type": "comparison",
                        "part": "content",
                        "title": "Compare",
                        "left": {"title": "Option A", "bullets": ["Fast"]},
                        "right": {"title": "Option B", "bullets": ["Safe"]},
                    },
                    {"type": "summary", "part": "closing", "title": "Summary", "key_points": ["One", "Two"], "next_steps": ["Act"]},
                    {"type": "table", "part": "content", "title": "Budget", "headers": ["Item", "Owner"], "rows": [["Design", "Alice"], ["Build", "Bob"]]},
                ],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            out_path = Path(temp_dir) / "demo.pptx"
            result = renderer.render(spec, out_path)
            self.assertTrue(result.path.exists())
            presentation = Presentation(str(result.path))
            self.assertEqual(len(presentation.slides), 6)

    def test_timeline_shapes_stay_within_slide_bounds(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "slides": [
                    {
                        "type": "timeline",
                        "part": "content",
                        "title": "Roadmap",
                        "events": [
                            {"label": "Q1", "title": "Plan", "detail": "Scope"},
                            {"label": "Q2", "title": "Build", "detail": "Execution"},
                            {"label": "Q3", "title": "Pilot", "detail": "Feedback"},
                            {"label": "Q4", "title": "Launch", "detail": "Scale"},
                        ],
                    }
                ],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            out_path = Path(temp_dir) / "timeline.pptx"
            result = renderer.render(spec, out_path)
            presentation = Presentation(str(result.path))
            slide = presentation.slides[0]
            slide_width = presentation.slide_width
            for shape in slide.shapes:
                self.assertGreaterEqual(int(shape.left), 0)
                self.assertLessEqual(int(shape.left + shape.width), slide_width)

    def test_graphic_slides_use_centered_safe_content_area(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "slides": [
                    {
                        "type": "image",
                        "part": "content",
                        "title": "Image",
                        "image": "placeholder",
                        "caption": "Local",
                        "bullets": ["Point"],
                    },
                    {
                        "type": "table",
                        "part": "content",
                        "title": "Budget",
                        "headers": ["Item", "Owner", "Status"],
                        "rows": [["Design", "Alice", "Done"], ["Build", "Bob", "Active"]],
                    },
                ],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            image_path = Path(temp_dir) / "pixel.png"
            image_path.write_bytes(PNG_BYTES)
            spec.slides[0].image = str(image_path)
            out_path = Path(temp_dir) / "graphics.pptx"
            result = renderer.render(spec, out_path)
            presentation = Presentation(str(result.path))
            slide_width = presentation.slide_width
            margin = Inches(0.85)

            image_slide = presentation.slides[0]
            pictures = [shape for shape in image_slide.shapes if shape.shape_type.name in {"PICTURE", "LINKED_PICTURE"}]
            self.assertEqual(len(pictures), 1)
            picture = pictures[0]
            self.assertGreaterEqual(int(picture.left), int(margin))
            self.assertLessEqual(int(picture.left + picture.width), slide_width - int(margin))

            table_slide = presentation.slides[1]
            table_shape = next(shape for shape in table_slide.shapes if getattr(shape, "has_table", False))
            self.assertGreaterEqual(int(table_shape.left), int(margin))
            self.assertLessEqual(int(table_shape.left + table_shape.width), slide_width - int(margin))

    def test_renderer_downgrades_missing_image_url_to_text_slide(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "slides": [
                    {
                        "type": "image",
                        "part": "content",
                        "title": "Image",
                        "image": "https://example.com/missing.png",
                        "caption": "Fallback",
                        "bullets": ["Note"],
                    }
                ],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            out_path = Path(temp_dir) / "image.pptx"
            with patch("office_agent.renderer.requests.get", side_effect=requests.RequestException("boom")):
                result = renderer.render(spec, out_path)
            self.assertTrue(result.warnings)
            presentation = Presentation(str(result.path))
            self.assertEqual(len(presentation.slides), 1)

    def test_renderer_embeds_local_image(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "slides": [
                    {
                        "type": "image",
                        "part": "content",
                        "title": "Image",
                        "image": "placeholder",
                        "caption": "Local",
                        "bullets": ["Point"],
                    }
                ],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            image_path = Path(temp_dir) / "pixel.png"
            image_path.write_bytes(PNG_BYTES)
            spec.slides[0].image = str(image_path)
            out_path = Path(temp_dir) / "local.pptx"
            result = renderer.render(spec, out_path)
            self.assertFalse(result.warnings)
            self.assertTrue(result.path.exists())

    def test_renderer_uses_template_slide_for_selected_part(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "template": {"opening": 0, "agenda": 1, "content": 1, "closing": 1},
                "slides": [{"type": "title", "part": "opening", "title": "Generated", "subtitle": "From template"}],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "template.pptx"
            template = Presentation()
            slide = template.slides.add_slide(template.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            title_box.name = "oa:title"
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
            subtitle_box.name = "oa:subtitle"
            template.save(str(template_path))

            out_path = Path(temp_dir) / "templated-output.pptx"
            result = renderer.render(spec, out_path, template_path=template_path)

            presentation = Presentation(str(result.path))
            self.assertEqual(len(presentation.slides), 1)
            texts = [shape.text for shape in presentation.slides[0].shapes if getattr(shape, "has_text_frame", False)]
            self.assertIn("Generated", texts)
            self.assertIn("From template", texts)

    def test_renderer_uses_native_powerpoint_placeholders_without_oa_markers(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "template": {"opening": 0, "agenda": 0, "content": 0, "closing": 0},
                "slides": [{"type": "bullets", "part": "content", "title": "Generated title", "bullets": ["Point A", "Point B"]}],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "native-template.pptx"
            template = Presentation()
            slide = template.slides.add_slide(template.slide_layouts[1])
            slide.shapes.title.text = "Old title"
            slide.placeholders[1].text = "Old body"
            template.save(str(template_path))

            out_path = Path(temp_dir) / "native-output.pptx"
            result = renderer.render(spec, out_path, template_path=template_path)

            self.assertFalse(result.warnings)
            presentation = Presentation(str(result.path))
            texts = [shape.text for shape in presentation.slides[0].shapes if getattr(shape, "has_text_frame", False)]
            joined = "\n".join(texts)
            self.assertIn("Generated title", joined)
            self.assertIn("Point A", joined)
            self.assertNotIn("Old body", joined)

    def test_renderer_preserves_branding_assets_while_removing_template_content(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "template": {"opening": 0, "agenda": 0, "content": 0, "closing": 0},
                "slides": [{"type": "bullets", "part": "content", "title": "Quarterly update", "bullets": ["Revenue up", "Hiring stable"]}],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "brand-template.pptx"
            logo_path = Path(temp_dir) / "logo.png"
            hero_path = Path(temp_dir) / "hero.png"
            logo_path.write_bytes(PNG_BYTES)
            hero_path.write_bytes(PNG_BYTES)

            template = Presentation()
            slide = template.slides.add_slide(template.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(6.5), Inches(0.8))
            title_box.name = "oa:title"
            title_box.text = "Old template title"
            body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(7.0), Inches(3.5))
            body_box.name = "oa:body"
            body_box.text = "Old template body"
            slide.shapes.add_picture(str(logo_path), Inches(11.2), Inches(0.3), width=Inches(1.1), height=Inches(0.5))
            slide.shapes.add_picture(str(hero_path), Inches(7.8), Inches(1.4), width=Inches(3.8), height=Inches(3.1))
            note_box = slide.shapes.add_textbox(Inches(8.0), Inches(5.3), Inches(3.2), Inches(0.8))
            note_box.text = "Template note to remove"
            template.save(str(template_path))

            out_path = Path(temp_dir) / "brand-output.pptx"
            result = renderer.render(spec, out_path, template_path=template_path)

            self.assertFalse(result.warnings)
            presentation = Presentation(str(result.path))
            rendered_slide = presentation.slides[0]
            texts = [shape.text for shape in rendered_slide.shapes if getattr(shape, "has_text_frame", False)]
            joined = "\n".join(texts)
            self.assertIn("Quarterly update", joined)
            self.assertIn("Revenue up", joined)
            self.assertNotIn("Template note to remove", joined)
            self.assertEqual(sum(1 for shape in rendered_slide.shapes if shape.shape_type.name in {"PICTURE", "LINKED_PICTURE"}), 1)

    def test_renderer_uses_blank_template_shell_for_content_pages(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "template": {"opening": 0, "agenda": 1, "content": 2, "closing": 3},
                "slides": [{"type": "bullets", "part": "content", "title": "Generated body", "bullets": ["Alpha", "Beta"]}],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "shell-template.pptx"
            template = Presentation()
            template.slides.add_slide(template.slide_layouts[6])
            template.slides.add_slide(template.slide_layouts[6])
            shell_slide = template.slides.add_slide(template.slide_layouts[6])
            shell_slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(1.2), Inches(0.12))
            template.slides.add_slide(template.slide_layouts[6])
            template.save(str(template_path))

            out_path = Path(temp_dir) / "shell-output.pptx"
            result = renderer.render(spec, out_path, template_path=template_path)

            self.assertFalse(result.warnings)
            presentation = Presentation(str(result.path))
            rendered_slide = presentation.slides[0]
            texts = [shape.text for shape in rendered_slide.shapes if getattr(shape, "has_text_frame", False)]
            self.assertIn("Generated body", "\n".join(texts))
            self.assertTrue(any(shape.shape_type.name == "AUTO_SHAPE" for shape in rendered_slide.shapes))

    def test_renderer_falls_back_when_template_placeholders_are_missing(self) -> None:
        spec = PresentationSpec.model_validate(
            {
                "title": "Demo",
                "template": {"opening": 0, "agenda": 0, "content": 0, "closing": 0},
                "slides": [{"type": "summary", "part": "closing", "title": "Wrap", "key_points": ["One"], "next_steps": ["Act"]}],
            }
        )
        renderer = PresentationRenderer()
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "template.pptx"
            template = Presentation()
            slide = template.slides.add_slide(template.slide_layouts[6])
            title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
            title_box.name = "oa:title"
            template.save(str(template_path))

            out_path = Path(temp_dir) / "fallback-output.pptx"
            result = renderer.render(spec, out_path, template_path=template_path)

            self.assertTrue(result.warnings)
            self.assertIn("does not expose writable roles", result.warnings[0])
            presentation = Presentation(str(result.path))
            self.assertEqual(len(presentation.slides), 1)
