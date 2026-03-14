from __future__ import annotations

import base64
import tempfile
from pathlib import Path
from unittest import TestCase
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches

from office_agent.template_preview import TemplatePreviewError, build_template_preview


PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9WnJyoAAAAAASUVORK5CYII="
)


class TemplatePreviewTests(TestCase):
    def test_build_template_preview_returns_slide_metadata(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "template.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
            title_shape.name = "oa:title"
            title_shape.text = "Cover"
            subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1))
            subtitle_shape.name = "oa:subtitle"
            presentation.save(str(template_path))

            def fake_export(_: Path, output_dir: Path) -> dict[int, Path]:
                image_path = output_dir / "slide-0.png"
                image_path.write_bytes(PNG_BYTES)
                return {0: image_path}

            with patch("office_agent.template_preview.export_template_thumbnails", side_effect=fake_export):
                preview = build_template_preview(template_path.read_bytes(), template_path.name)

        self.assertEqual(len(preview.slides), 1)
        self.assertEqual(preview.slides[0].title_text, "Cover")
        self.assertEqual(preview.slides[0].placeholder_roles, ["subtitle", "title"])
        self.assertTrue(preview.slides[0].thumbnail_data_url.startswith("data:image/png;base64,"))
        self.assertEqual(preview.cleanup_mode, "preserve_branding")

    def test_build_template_preview_detects_native_placeholders(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "native-template.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "Agenda"
            slide.placeholders[1].text = "Body placeholder"
            presentation.save(str(template_path))

            def fake_export(_: Path, output_dir: Path) -> dict[int, Path]:
                image_path = output_dir / "slide-0.png"
                image_path.write_bytes(PNG_BYTES)
                return {0: image_path}

            with patch("office_agent.template_preview.export_template_thumbnails", side_effect=fake_export):
                preview = build_template_preview(template_path.read_bytes(), template_path.name)

        self.assertEqual(preview.slides[0].placeholder_roles, ["body", "title"])

    def test_build_template_preview_rejects_non_pptx(self) -> None:
        with self.assertRaises(TemplatePreviewError):
            build_template_preview(b"not-a-pptx", "template.txt")
