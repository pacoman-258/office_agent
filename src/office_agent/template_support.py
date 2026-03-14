from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass

from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE


PLACEHOLDER_PREFIX = "oa:"
SUPPORTED_PLACEHOLDER_ROLES = {"title", "subtitle", "body", "image", "caption"}
TEMPLATE_CLEANUP_MODE = "preserve_branding"

TITLE_PLACEHOLDER_TYPES = {PP_PLACEHOLDER_TYPE.TITLE, PP_PLACEHOLDER_TYPE.CENTER_TITLE, PP_PLACEHOLDER_TYPE.VERTICAL_TITLE}
SUBTITLE_PLACEHOLDER_TYPES = {PP_PLACEHOLDER_TYPE.SUBTITLE}
BODY_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER_TYPE.BODY,
    PP_PLACEHOLDER_TYPE.OBJECT,
    PP_PLACEHOLDER_TYPE.VERTICAL_BODY,
    PP_PLACEHOLDER_TYPE.VERTICAL_OBJECT,
}
IMAGE_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER_TYPE.PICTURE,
    PP_PLACEHOLDER_TYPE.BITMAP,
    PP_PLACEHOLDER_TYPE.MEDIA_CLIP,
    PP_PLACEHOLDER_TYPE.SLIDE_IMAGE,
}
DECORATIVE_PLACEHOLDER_TYPES = {
    PP_PLACEHOLDER_TYPE.DATE,
    PP_PLACEHOLDER_TYPE.FOOTER,
    PP_PLACEHOLDER_TYPE.HEADER,
    PP_PLACEHOLDER_TYPE.SLIDE_NUMBER,
}
CONTENT_SHAPE_TYPES = {
    MSO_SHAPE_TYPE.CHART,
    MSO_SHAPE_TYPE.DIAGRAM,
    MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
    MSO_SHAPE_TYPE.IGX_GRAPHIC,
    MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
    MSO_SHAPE_TYPE.MEDIA,
    MSO_SHAPE_TYPE.OLE_CONTROL_OBJECT,
    MSO_SHAPE_TYPE.TABLE,
    MSO_SHAPE_TYPE.WEB_VIDEO,
}
PICTURE_SHAPE_TYPES = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}


@dataclass(frozen=True)
class TemplateSlot:
    role: str
    shape: object | None
    left: int
    top: int
    width: int
    height: int
    source: str


@dataclass(frozen=True)
class TemplateShapeAction:
    clear_shape_ids: set[int]
    delete_shape_ids: set[int]


@dataclass(frozen=True)
class TemplateAnalysis:
    slots: dict[str, TemplateSlot]
    actions: TemplateShapeAction

    @property
    def placeholder_roles(self) -> list[str]:
        return sorted(self.slots)

    @property
    def cleanup_mode(self) -> str:
        return TEMPLATE_CLEANUP_MODE


def duplicate_slide(prs, source_slide):
    layout = source_slide.slide_layout
    duplicated = prs.slides.add_slide(layout)

    for shape in list(duplicated.shapes):
        duplicated.shapes._spTree.remove(shape.element)

    for shape in source_slide.shapes:
        duplicated.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

    for rel in source_slide.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        if rel.is_external:
            duplicated.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            duplicated.part.rels.get_or_add(rel.reltype, rel.target_part)

    return duplicated


def analyze_template_slide(slide) -> TemplateAnalysis:
    slots = _detect_slots(slide)
    actions = _plan_cleanup(slide, slots)
    return TemplateAnalysis(slots=slots, actions=actions)


def extract_placeholder_roles(slide) -> list[str]:
    return analyze_template_slide(slide).placeholder_roles


def extract_shape_role(shape) -> str | None:
    for candidate in _shape_markers(shape):
        normalized = candidate.strip().lower()
        if not normalized.startswith(PLACEHOLDER_PREFIX):
            continue
        role = normalized.removeprefix(PLACEHOLDER_PREFIX)
        if role in SUPPORTED_PLACEHOLDER_ROLES:
            return role
    return None


def extract_slide_title_text(slide) -> str | None:
    analysis = analyze_template_slide(slide)
    title_slot = analysis.slots.get("title")
    if title_slot is not None and title_slot.shape is not None and getattr(title_slot.shape, "has_text_frame", False):
        text = title_slot.shape.text.strip()
        if text:
            return text

    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text.strip()
            if text:
                return text
    return None


def sanitize_template_slide(slide, analysis: TemplateAnalysis) -> None:
    for shape in list(slide.shapes):
        shape_id = getattr(shape, "shape_id", None)
        if shape_id in analysis.actions.delete_shape_ids:
            _remove_shape(shape)
            continue
        if shape_id in analysis.actions.clear_shape_ids:
            _clear_shape(shape)


def is_template_shell(slide, analysis: TemplateAnalysis | None = None) -> bool:
    analysis = analysis or analyze_template_slide(slide)
    if analysis.placeholder_roles:
        return False
    for shape in slide.shapes:
        if _should_delete_shape(slide, shape):
            return False
        if getattr(shape, "has_text_frame", False) and shape.text.strip() and not _is_branding_text(slide, shape):
            return False
        if shape.shape_type in CONTENT_SHAPE_TYPES:
            return False
        if shape.shape_type in PICTURE_SHAPE_TYPES and not _looks_like_branding_asset(slide, shape):
            return False
    return True


def _detect_slots(slide) -> dict[str, TemplateSlot]:
    slots: dict[str, TemplateSlot] = {}
    for shape in slide.shapes:
        role = extract_shape_role(shape)
        if role and role not in slots:
            slots[role] = _slot_from_shape(role, shape, "marker")

    shapes = list(slide.shapes)
    text_candidates = [shape for shape in shapes if _is_text_candidate(shape)]
    image_candidates = [shape for shape in shapes if _is_image_candidate(shape)]

    _assign_slot_from_candidates(slots, "title", _title_candidates(slide, text_candidates), "native")
    _assign_slot_from_candidates(slots, "subtitle", _subtitle_candidates(slide, text_candidates, slots.get("title")), "native")
    _assign_slot_from_candidates(slots, "body", _body_candidates(slide, text_candidates, slots), "native")
    _assign_slot_from_candidates(slots, "image", _image_candidates(slide, image_candidates, slots), "native")
    _assign_slot_from_candidates(slots, "caption", _caption_candidates(slide, text_candidates, slots), "heuristic")
    return slots


def _assign_slot_from_candidates(
    slots: dict[str, TemplateSlot],
    role: str,
    candidates: list[object],
    source: str,
) -> None:
    if role in slots:
        return
    for shape in candidates:
        if _shape_id(shape) in {_shape_id(slot.shape) for slot in slots.values() if slot.shape is not None}:
            continue
        slots[role] = _slot_from_shape(role, shape, source)
        return


def _slot_from_shape(role: str, shape, source: str) -> TemplateSlot:
    return TemplateSlot(
        role=role,
        shape=shape,
        left=int(shape.left),
        top=int(shape.top),
        width=int(shape.width),
        height=int(shape.height),
        source=source,
    )


def _plan_cleanup(slide, slots: dict[str, TemplateSlot]) -> TemplateShapeAction:
    clear_shape_ids: set[int] = set()
    delete_shape_ids: set[int] = set()
    selected_slots = {_shape_id(slot.shape): slot for slot in slots.values() if slot.shape is not None}

    for shape in slide.shapes:
        shape_id = _shape_id(shape)
        if shape_id in selected_slots:
            role = selected_slots[shape_id].role
            if role == "image" and _is_replaced_picture(shape):
                delete_shape_ids.add(shape_id)
            elif role in {"title", "subtitle", "body", "caption"} and getattr(shape, "has_text_frame", False):
                clear_shape_ids.add(shape_id)
            continue

        if _should_delete_shape(slide, shape):
            delete_shape_ids.add(shape_id)

    return TemplateShapeAction(clear_shape_ids=clear_shape_ids, delete_shape_ids=delete_shape_ids)


def _should_delete_shape(slide, shape) -> bool:
    placeholder_type = _placeholder_type(shape)
    if placeholder_type in DECORATIVE_PLACEHOLDER_TYPES:
        return False
    if extract_shape_role(shape):
        return False
    if shape.shape_type in CONTENT_SHAPE_TYPES:
        return True
    if shape.shape_type in PICTURE_SHAPE_TYPES:
        return not _looks_like_branding_asset(slide, shape)
    if getattr(shape, "has_table", False) or getattr(shape, "has_chart", False):
        return True
    if getattr(shape, "has_text_frame", False):
        if not shape.text.strip():
            return False
        if _is_branding_text(slide, shape):
            return False
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return placeholder_type not in DECORATIVE_PLACEHOLDER_TYPES
    return False


def _clear_shape(shape) -> None:
    if getattr(shape, "has_text_frame", False):
        frame = shape.text_frame
        frame.clear()
        if frame.paragraphs:
            frame.paragraphs[0].text = ""


def _remove_shape(shape) -> None:
    element = getattr(shape, "_element", None)
    parent = getattr(element, "getparent", lambda: None)()
    if element is not None and parent is not None:
        parent.remove(element)


def _is_text_candidate(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    placeholder_type = _placeholder_type(shape)
    if placeholder_type in DECORATIVE_PLACEHOLDER_TYPES:
        return False
    text = shape.text.strip()
    return bool(text) or placeholder_type in TITLE_PLACEHOLDER_TYPES | SUBTITLE_PLACEHOLDER_TYPES | BODY_PLACEHOLDER_TYPES


def _is_image_candidate(shape) -> bool:
    placeholder_type = _placeholder_type(shape)
    if placeholder_type in IMAGE_PLACEHOLDER_TYPES | BODY_PLACEHOLDER_TYPES:
        return True
    return shape.shape_type in PICTURE_SHAPE_TYPES


def _title_candidates(slide, candidates: list[object]) -> list[object]:
    direct_title = getattr(slide.shapes, "title", None)
    slide_width, slide_height = _slide_dimensions(slide)
    ranked: list[tuple[float, object]] = []
    if direct_title is not None and _is_text_candidate(direct_title):
        ranked.append((1000.0, direct_title))

    for shape in candidates:
        placeholder_type = _placeholder_type(shape)
        top = int(shape.top)
        height = max(int(shape.height), 1)
        area = int(shape.width) * height
        score = float(area) / max(top + 1, 1)
        if placeholder_type in TITLE_PLACEHOLDER_TYPES:
            score += 500.0
        if top <= slide_height * 0.35:
            score += 200.0
        ranked.append((score, shape))
    return [shape for _, shape in sorted(ranked, key=lambda item: item[0], reverse=True)]


def _subtitle_candidates(slide, candidates: list[object], title_slot: TemplateSlot | None) -> list[object]:
    title_shape_id = _shape_id(title_slot.shape) if title_slot is not None else None
    slide_width, slide_height = _slide_dimensions(slide)
    ranked: list[tuple[float, object]] = []
    for shape in candidates:
        if _shape_id(shape) == title_shape_id:
            continue
        placeholder_type = _placeholder_type(shape)
        if placeholder_type in BODY_PLACEHOLDER_TYPES | IMAGE_PLACEHOLDER_TYPES:
            continue
        top = int(shape.top)
        area = int(shape.width) * max(int(shape.height), 1)
        score = float(area) / max(top + 1, 1)
        if placeholder_type in SUBTITLE_PLACEHOLDER_TYPES:
            score += 500.0
        elif not shape.text.strip():
            continue
        elif area > slide_width * slide_height * 0.12:
            continue
        if title_slot is not None and top >= title_slot.top:
            score += 150.0
        if top <= slide_height * 0.55:
            score += 80.0
        ranked.append((score, shape))
    return [shape for _, shape in sorted(ranked, key=lambda item: item[0], reverse=True)]


def _body_candidates(slide, candidates: list[object], slots: dict[str, TemplateSlot]) -> list[object]:
    reserved = {_shape_id(slot.shape) for slot in slots.values() if slot.shape is not None}
    _, slide_height = _slide_dimensions(slide)
    ranked: list[tuple[float, object]] = []
    for shape in candidates:
        if _shape_id(shape) in reserved:
            continue
        placeholder_type = _placeholder_type(shape)
        area = int(shape.width) * max(int(shape.height), 1)
        score = float(area)
        if placeholder_type in BODY_PLACEHOLDER_TYPES:
            score += 800.0
        if int(shape.top) >= slide_height * 0.18:
            score += 120.0
        if not _is_branding_text(slide, shape):
            ranked.append((score, shape))
    return [shape for _, shape in sorted(ranked, key=lambda item: item[0], reverse=True)]


def _image_candidates(slide, candidates: list[object], slots: dict[str, TemplateSlot]) -> list[object]:
    reserved = {_shape_id(slot.shape) for slot in slots.values() if slot.shape is not None}
    ranked: list[tuple[float, object]] = []
    for shape in candidates:
        if _shape_id(shape) in reserved:
            continue
        placeholder_type = _placeholder_type(shape)
        area = int(shape.width) * max(int(shape.height), 1)
        score = float(area)
        if placeholder_type in IMAGE_PLACEHOLDER_TYPES:
            score += 700.0
        if shape.shape_type in PICTURE_SHAPE_TYPES and not _looks_like_branding_asset(slide, shape):
            score += 400.0
        if area > 0 and not _looks_like_branding_asset(slide, shape):
            ranked.append((score, shape))
    return [shape for _, shape in sorted(ranked, key=lambda item: item[0], reverse=True)]


def _caption_candidates(slide, candidates: list[object], slots: dict[str, TemplateSlot]) -> list[object]:
    reserved = {_shape_id(slot.shape) for slot in slots.values() if slot.shape is not None}
    image_slot = slots.get("image")
    slide_width, slide_height = _slide_dimensions(slide)
    ranked: list[tuple[float, object]] = []
    for shape in candidates:
        if _shape_id(shape) in reserved:
            continue
        area = int(shape.width) * max(int(shape.height), 1)
        if area > slide_width * slide_height * 0.12:
            continue
        score = float(area)
        if image_slot is not None and int(shape.top) >= image_slot.top:
            score += 150.0
        if _is_branding_text(slide, shape):
            continue
        ranked.append((score, shape))
    return [shape for _, shape in sorted(ranked, key=lambda item: item[0], reverse=True)]


def _placeholder_type(shape):
    try:
        if not getattr(shape, "is_placeholder", False):
            return None
        return shape.placeholder_format.type
    except Exception:
        return None


def _shape_markers(shape) -> list[str]:
    markers: list[str] = []
    name = getattr(shape, "name", None)
    if isinstance(name, str) and name.strip():
        markers.append(name)

    c_nv_pr = getattr(getattr(shape, "_element", None), "_nvXxPr", None)
    c_nv_pr = getattr(c_nv_pr, "cNvPr", None)
    if c_nv_pr is not None:
        for attribute in ("descr", "title"):
            value = c_nv_pr.get(attribute)
            if value:
                markers.append(value)

    return markers


def _shape_id(shape) -> int | None:
    return getattr(shape, "shape_id", None) if shape is not None else None


def _shape_area(slide, shape) -> float:
    slide_width, slide_height = _slide_dimensions(slide)
    return (int(shape.width) * int(shape.height)) / float(slide_width * slide_height)


def _looks_like_branding_asset(slide, shape) -> bool:
    slide_width, slide_height = _slide_dimensions(slide)
    area_ratio = _shape_area(slide, shape)
    left = int(shape.left)
    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)
    right_gap = int(slide_width) - (left + width)
    bottom_gap = int(slide_height) - (top + height)
    near_edge = top < slide_height * 0.16 or left < slide_width * 0.12 or right_gap < slide_width * 0.12 or bottom_gap < slide_height * 0.14
    return area_ratio <= 0.05 and near_edge


def _is_branding_text(slide, shape) -> bool:
    placeholder_type = _placeholder_type(shape)
    if placeholder_type in DECORATIVE_PLACEHOLDER_TYPES:
        return True
    text = shape.text.strip()
    if not text:
        return False
    return _looks_like_branding_asset(slide, shape) and len(text) <= 40


def _is_replaced_picture(shape) -> bool:
    return shape.shape_type in PICTURE_SHAPE_TYPES


def _slide_dimensions(slide) -> tuple[int, int]:
    presentation = slide.part.package.presentation_part.presentation
    sld_sz = presentation._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}sldSz")
    if sld_sz is None:
        raise ValueError("Unable to determine slide dimensions from presentation XML.")
    return int(sld_sz.get("cx")), int(sld_sz.get("cy"))
