from __future__ import annotations

import os
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.parts.image import Image as PptxImage
from pptx.util import Inches, Pt

from office_agent.errors import RenderError
from office_agent.schema import (
    BulletsSlideSpec,
    ComparisonSlideSpec,
    ImageSlideSpec,
    PresentationSpec,
    QuoteSlideSpec,
    SectionSlideSpec,
    SlideSpec,
    SummarySlideSpec,
    TableSlideSpec,
    TemplatePart,
    TemplateSelectionSpec,
    ThemeSpec,
    TimelineSlideSpec,
    TitleSlideSpec,
    TwoColumnSlideSpec,
)
from office_agent.template_support import duplicate_slide, extract_placeholder_roles, extract_shape_role


@dataclass(frozen=True)
class ThemeTokens:
    primary_color: str
    accent_color: str
    background_color: str
    body_color: str
    muted_color: str
    heading_font: str
    body_font: str
    cover_style: str


PRESET_THEMES: dict[str, ThemeTokens] = {
    "default": ThemeTokens(
        primary_color="#16324F",
        accent_color="#2A6F97",
        background_color="#F5F7FA",
        body_color="#1F2937",
        muted_color="#6B7280",
        heading_font="Microsoft YaHei",
        body_font="Microsoft YaHei",
        cover_style="band",
    ),
    "executive": ThemeTokens(
        primary_color="#0E2A47",
        accent_color="#C28B2C",
        background_color="#F7F3EC",
        body_color="#243647",
        muted_color="#6B6F75",
        heading_font="Georgia",
        body_font="Microsoft YaHei",
        cover_style="centered",
    ),
    "editorial": ThemeTokens(
        primary_color="#3A243B",
        accent_color="#D95D39",
        background_color="#FFF9F1",
        body_color="#403332",
        muted_color="#7A6A68",
        heading_font="Georgia",
        body_font="Microsoft YaHei",
        cover_style="minimal",
    ),
}


@dataclass
class RenderResult:
    path: Path
    warnings: list[str] = field(default_factory=list)


class PresentationRenderer:
    def render(
        self,
        spec: PresentationSpec,
        out_path: str | Path,
        template_path: str | Path | None = None,
    ) -> RenderResult:
        output_path = Path(out_path)
        prs = Presentation(str(template_path)) if template_path else Presentation()
        if template_path is None:
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        tokens = resolve_theme_tokens(spec.theme)
        warnings: list[str] = []
        temp_files: list[Path] = []
        template_source_slides = list(prs.slides) if template_path and spec.template else []
        template_source_count = len(template_source_slides)

        try:
            for slide_spec in spec.slides:
                if not self._render_with_template(prs, slide_spec, spec.template, template_source_slides, tokens, warnings, temp_files):
                    self._render_default_slide(prs, slide_spec, tokens, warnings, temp_files)

            if template_source_count:
                self._remove_template_source_slides(prs, template_source_count)

            output_path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(str(output_path))
            return RenderResult(path=output_path, warnings=warnings)
        except Exception as exc:  # pragma: no cover
            raise RenderError(f"Failed to render presentation: {exc}") from exc
        finally:
            for temp_path in temp_files:
                try:
                    temp_path.unlink(missing_ok=True)
                except OSError:
                    pass

    def _render_with_template(
        self,
        prs: Presentation,
        slide_spec: SlideSpec,
        template: TemplateSelectionSpec | None,
        template_source_slides: list,
        tokens: ThemeTokens,
        warnings: list[str],
        temp_files: list[Path],
    ) -> bool:
        if template is None or not template_source_slides:
            return False

        source_index = getattr(template, slide_spec.part)
        if source_index >= len(template_source_slides):
            warnings.append(
                f"Template slide index {source_index} for '{slide_spec.part}' is out of range. Falling back to default rendering."
            )
            return False

        source_slide = template_source_slides[source_index]
        available_roles = set(extract_placeholder_roles(source_slide))
        required_roles = required_template_roles(slide_spec)
        if not required_roles.issubset(available_roles):
            warnings.append(
                f"Template slide {source_index} for '{slide_spec.part}' is missing placeholders {sorted(required_roles - available_roles)}. Falling back to default rendering."
            )
            return False

        if isinstance(slide_spec, ImageSlideSpec):
            image_path, _, cleanup = self._resolve_image_source(slide_spec.image)
            if image_path is None:
                if cleanup is not None:
                    temp_files.append(cleanup)
                return False
            if cleanup is not None:
                temp_files.append(cleanup)
        else:
            image_path = None

        page = duplicate_slide(prs, source_slide)
        placeholder_map = self._map_template_shapes(page)
        self._fill_text_placeholder(placeholder_map.get("title"), slide_spec.title)
        self._fill_text_placeholder(placeholder_map.get("subtitle"), self._subtitle_text(slide_spec))
        self._fill_body_placeholder(placeholder_map.get("body"), body_lines_for_slide(slide_spec))
        self._fill_text_placeholder(placeholder_map.get("caption"), slide_spec.caption if isinstance(slide_spec, ImageSlideSpec) else None)
        if isinstance(slide_spec, ImageSlideSpec) and image_path is not None:
            self._fill_image_placeholder(page, placeholder_map.get("image"), image_path)
        return True

    def _render_default_slide(
        self,
        prs: Presentation,
        slide_spec: SlideSpec,
        tokens: ThemeTokens,
        warnings: list[str],
        temp_files: list[Path],
    ) -> None:
        if isinstance(slide_spec, TitleSlideSpec):
            self._render_title_slide(prs, slide_spec, tokens)
        elif isinstance(slide_spec, SectionSlideSpec):
            self._render_section_slide(prs, slide_spec, tokens)
        elif isinstance(slide_spec, BulletsSlideSpec):
            self._render_bullets_slide(prs, slide_spec, tokens)
        elif isinstance(slide_spec, TwoColumnSlideSpec):
            self._render_two_column_slide(prs, slide_spec, tokens)
        elif isinstance(slide_spec, ImageSlideSpec):
            temp_path, warning = self._render_image_slide(prs, slide_spec, tokens)
            if temp_path is not None:
                temp_files.append(temp_path)
            if warning:
                warnings.append(warning)
        elif isinstance(slide_spec, TimelineSlideSpec):
            self._render_timeline_slide(prs, slide_spec, tokens)
        elif isinstance(slide_spec, QuoteSlideSpec):
            self._render_quote_slide(prs, slide_spec, tokens)
        elif isinstance(slide_spec, ComparisonSlideSpec):
            self._render_comparison_slide(prs, slide_spec, tokens)
        elif isinstance(slide_spec, SummarySlideSpec):
            self._render_summary_slide(prs, slide_spec, tokens)
        elif isinstance(slide_spec, TableSlideSpec):
            self._render_table_slide(prs, slide_spec, tokens)
        else:
            raise RenderError(f"Unsupported slide spec: {slide_spec!r}")

    def _render_title_slide(self, prs: Presentation, slide: TitleSlideSpec, tokens: ThemeTokens) -> None:
        page = self._new_slide(prs, tokens)
        if tokens.cover_style == "band":
            band = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
            self._fill_shape(band, tokens.primary_color, tokens.primary_color)

        title_box = page.shapes.add_textbox(Inches(0.95), Inches(1.45), Inches(11.2), Inches(1.7))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = slide.title
        title_p.font.name = tokens.heading_font
        title_p.font.size = Pt(28 if tokens.cover_style != "minimal" else 26)
        title_p.font.bold = True
        title_p.font.color.rgb = rgb(tokens.primary_color)
        title_p.alignment = PP_ALIGN.CENTER if tokens.cover_style == "centered" else PP_ALIGN.LEFT

        if slide.subtitle:
            subtitle_box = page.shapes.add_textbox(Inches(0.95), Inches(3.0), Inches(10.9), Inches(1.3))
            subtitle_p = subtitle_box.text_frame.paragraphs[0]
            subtitle_p.text = slide.subtitle
            subtitle_p.font.name = tokens.body_font
            subtitle_p.font.size = Pt(18)
            subtitle_p.font.color.rgb = rgb(tokens.muted_color)
            subtitle_p.alignment = title_p.alignment

        if tokens.cover_style != "minimal":
            line = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), Inches(4.3), Inches(3.4), Inches(0.06))
            self._fill_shape(line, tokens.accent_color, tokens.accent_color)

    def _render_section_slide(self, prs: Presentation, slide: SectionSlideSpec, tokens: ThemeTokens) -> None:
        page = self._new_slide(prs, tokens)
        badge = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.4), Inches(1.7), Inches(0.42))
        self._fill_shape(badge, tokens.accent_color, tokens.accent_color)
        badge.text = "Section"
        self._style_shape_text(badge, tokens.body_font, 11, "#FFFFFF", bold=True, align=PP_ALIGN.CENTER)

        title_box = page.shapes.add_textbox(Inches(0.95), Inches(2.2), Inches(11.0), Inches(1.5))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = slide.title
        title_p.font.name = tokens.heading_font
        title_p.font.size = Pt(30)
        title_p.font.bold = True
        title_p.font.color.rgb = rgb(tokens.primary_color)

        if slide.subtitle:
            subtitle_box = page.shapes.add_textbox(Inches(0.95), Inches(3.8), Inches(10.5), Inches(1.0))
            subtitle_p = subtitle_box.text_frame.paragraphs[0]
            subtitle_p.text = slide.subtitle
            subtitle_p.font.name = tokens.body_font
            subtitle_p.font.size = Pt(16)
            subtitle_p.font.color.rgb = rgb(tokens.muted_color)

    def _render_bullets_slide(self, prs: Presentation, slide: BulletsSlideSpec, tokens: ThemeTokens) -> None:
        page = self._new_slide(prs, tokens)
        self._add_slide_title(page, slide.title, tokens)
        textbox = page.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.2), Inches(4.9))
        self._write_bullets(textbox.text_frame, slide.bullets, tokens)

    def _render_two_column_slide(self, prs: Presentation, slide: TwoColumnSlideSpec, tokens: ThemeTokens) -> None:
        page = self._new_slide(prs, tokens)
        self._add_slide_title(page, slide.title, tokens)
        left_box = page.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(5.2), Inches(4.8))
        right_box = page.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.2), Inches(4.8))
        self._write_column(left_box.text_frame, slide.left_title, slide.left_bullets, tokens)
        self._write_column(right_box.text_frame, slide.right_title, slide.right_bullets, tokens)
        divider = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.48), Inches(1.85), Inches(0.03), Inches(4.6))
        self._fill_shape(divider, tokens.accent_color, tokens.accent_color)

    def _render_image_slide(self, prs: Presentation, slide: ImageSlideSpec, tokens: ThemeTokens) -> tuple[Path | None, str | None]:
        page = self._new_slide(prs, tokens)
        self._add_slide_title(page, slide.title, tokens)
        image_path, warning, cleanup = self._resolve_image_source(slide.image)
        if image_path:
            self._add_image(page, image_path)
            if slide.caption:
                caption_box = page.shapes.add_textbox(Inches(0.9), Inches(6.25), Inches(8.0), Inches(0.45))
                caption_p = caption_box.text_frame.paragraphs[0]
                caption_p.text = slide.caption
                caption_p.font.name = tokens.body_font
                caption_p.font.size = Pt(11)
                caption_p.font.color.rgb = rgb(tokens.muted_color)
        else:
            placeholder_box = page.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(7.8), Inches(2.0))
            placeholder_p = placeholder_box.text_frame.paragraphs[0]
            placeholder_p.text = "Image unavailable. This slide was downgraded to a text summary."
            placeholder_p.font.name = tokens.body_font
            placeholder_p.font.size = Pt(20)
            placeholder_p.font.bold = True
            placeholder_p.font.color.rgb = rgb(tokens.accent_color)

        notes = slide.bullets or ([slide.caption] if slide.caption else ["Add supporting notes here."])
        bullet_box = page.shapes.add_textbox(Inches(9.2), Inches(1.8), Inches(3.1), Inches(4.6))
        self._write_bullets(bullet_box.text_frame, notes, tokens)
        return cleanup, warning

    def _render_timeline_slide(self, prs: Presentation, slide: TimelineSlideSpec, tokens: ThemeTokens) -> None:
        page = self._new_slide(prs, tokens)
        self._add_slide_title(page, slide.title, tokens)
        line = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(3.6), Inches(11.0), Inches(0.05))
        self._fill_shape(line, tokens.accent_color, tokens.accent_color)
        step_gap = 10.8 / max(1, len(slide.events) - 1)
        for index, event in enumerate(slide.events):
            x = 1.0 + index * step_gap
            dot = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(3.4), Inches(0.22), Inches(0.22))
            self._fill_shape(dot, tokens.accent_color, tokens.accent_color)
            label_box = page.shapes.add_textbox(Inches(x - 0.25), Inches(2.75), Inches(1.0), Inches(0.4))
            label_p = label_box.text_frame.paragraphs[0]
            label_p.text = event.label
            label_p.font.name = tokens.body_font
            label_p.font.size = Pt(11)
            label_p.font.bold = True
            label_p.font.color.rgb = rgb(tokens.accent_color)

            title_box = page.shapes.add_textbox(Inches(x - 0.25), Inches(3.8), Inches(1.7), Inches(1.2))
            frame = title_box.text_frame
            frame.word_wrap = True
            p = frame.paragraphs[0]
            p.text = event.title
            p.font.name = tokens.heading_font
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = rgb(tokens.primary_color)
            if event.detail:
                detail = frame.add_paragraph()
                detail.text = event.detail
                detail.font.name = tokens.body_font
                detail.font.size = Pt(11)
                detail.font.color.rgb = rgb(tokens.muted_color)

    def _render_quote_slide(self, prs: Presentation, slide: QuoteSlideSpec, tokens: ThemeTokens) -> None:
        page = self._new_slide(prs, tokens)
        quote_block = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.2), Inches(11.4), Inches(4.6))
        self._fill_shape(quote_block, lighten(tokens.background_color), tokens.accent_color)
        text_frame = quote_block.text_frame
        text_frame.margin_left = Inches(0.35)
        text_frame.margin_right = Inches(0.35)
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        quote_p = text_frame.paragraphs[0]
        quote_p.text = f'"{slide.quote}"'
        quote_p.font.name = tokens.heading_font
        quote_p.font.size = Pt(24)
        quote_p.font.bold = True
        quote_p.font.color.rgb = rgb(tokens.primary_color)
        quote_p.alignment = PP_ALIGN.CENTER

        if slide.attribution or slide.source:
            byline = text_frame.add_paragraph()
            byline.text = " - ".join(part for part in [slide.attribution, slide.source] if part)
            byline.font.name = tokens.body_font
            byline.font.size = Pt(14)
            byline.font.color.rgb = rgb(tokens.muted_color)
            byline.alignment = PP_ALIGN.CENTER

        self._add_slide_title(page, slide.title, tokens)

    def _render_comparison_slide(self, prs: Presentation, slide: ComparisonSlideSpec, tokens: ThemeTokens) -> None:
        page = self._new_slide(prs, tokens)
        self._add_slide_title(page, slide.title, tokens)
        self._render_panel_with_bullets(page, Inches(0.9), slide.left.title, slide.left.bullets, tokens)
        self._render_panel_with_bullets(page, Inches(6.7), slide.right.title, slide.right.bullets, tokens)

    def _render_summary_slide(self, prs: Presentation, slide: SummarySlideSpec, tokens: ThemeTokens) -> None:
        page = self._new_slide(prs, tokens)
        self._add_slide_title(page, slide.title, tokens)

        summary_panel = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.7), Inches(6.0), Inches(4.9))
        self._fill_shape(summary_panel, lighten(tokens.background_color), tokens.accent_color)
        self._style_shape_text(summary_panel, tokens.heading_font, 16, tokens.primary_color, bold=True)
        summary_panel.text = "Key Points"

        summary_text = page.shapes.add_textbox(Inches(1.2), Inches(2.25), Inches(5.3), Inches(3.8))
        self._write_bullets(summary_text.text_frame, slide.key_points, tokens)

        next_steps_panel = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(1.7), Inches(5.0), Inches(4.9))
        self._fill_shape(next_steps_panel, tokens.primary_color, tokens.primary_color)
        self._style_shape_text(next_steps_panel, tokens.heading_font, 16, "#FFFFFF", bold=True)
        next_steps_panel.text = "Next Steps"

        next_steps_text = page.shapes.add_textbox(Inches(7.6), Inches(2.25), Inches(4.3), Inches(3.8))
        self._write_bullets(next_steps_text.text_frame, slide.next_steps or ["No explicit next steps provided."], tokens, font_color="#FFFFFF")

    def _render_table_slide(self, prs: Presentation, slide: TableSlideSpec, tokens: ThemeTokens) -> None:
        page = self._new_slide(prs, tokens)
        self._add_slide_title(page, slide.title, tokens)
        rows = len(slide.rows) + 1
        cols = len(slide.headers)
        table = page.shapes.add_table(rows, cols, Inches(0.9), Inches(1.7), Inches(11.5), Inches(4.9)).table

        for col_idx, header in enumerate(slide.headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            self._style_table_cell(cell, tokens.heading_font, 13, "#FFFFFF", tokens.primary_color)

        for row_idx, row in enumerate(slide.rows, start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                fill_color = lighten(tokens.background_color) if row_idx % 2 == 1 else "#FFFFFF"
                self._style_table_cell(cell, tokens.body_font, 11, tokens.body_color, fill_color)

    def _new_slide(self, prs: Presentation, tokens: ThemeTokens):
        layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[len(prs.slide_layouts) - 1]
        page = prs.slides.add_slide(layout)
        page.background.fill.solid()
        page.background.fill.fore_color.rgb = rgb(tokens.background_color)
        return page

    def _map_template_shapes(self, slide) -> dict[str, object]:
        return {role: shape for shape in slide.shapes if (role := extract_shape_role(shape))}

    def _fill_text_placeholder(self, shape, value: str | None) -> None:
        if shape is None or value is None or not getattr(shape, "has_text_frame", False):
            return
        shape.text = value

    def _fill_body_placeholder(self, shape, lines: list[str]) -> None:
        if shape is None or not lines or not getattr(shape, "has_text_frame", False):
            return
        frame = shape.text_frame
        frame.clear()
        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0

    def _fill_image_placeholder(self, slide, shape, image_path: Path) -> None:
        if shape is None:
            return
        slide.shapes.add_picture(str(image_path), shape.left, shape.top, width=shape.width, height=shape.height)

    def _remove_template_source_slides(self, prs: Presentation, count: int) -> None:
        for index in range(count - 1, -1, -1):
            slide_id = prs.slides._sldIdLst[index]
            rel_id = slide_id.rId
            prs.part.drop_rel(rel_id)
            prs.slides._sldIdLst.remove(slide_id)

    def _render_panel_with_bullets(self, page, left: int, title: str, bullets: list[str], tokens: ThemeTokens) -> None:
        panel = page.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.8), Inches(5.7), Inches(4.9))
        self._fill_shape(panel, lighten(tokens.background_color), tokens.accent_color)
        self._style_shape_text(panel, tokens.heading_font, 16, tokens.primary_color, bold=True)
        panel.text = title
        body_box = page.shapes.add_textbox(left + Inches(0.28), Inches(2.35), Inches(5.1), Inches(3.9))
        self._write_bullets(body_box.text_frame, bullets, tokens)

    def _add_slide_title(self, slide, title: str, tokens: ThemeTokens) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.45), Inches(11.4), Inches(0.8))
        title_p = title_box.text_frame.paragraphs[0]
        title_p.text = title
        title_p.font.name = tokens.heading_font
        title_p.font.size = Pt(22)
        title_p.font.bold = True
        title_p.font.color.rgb = rgb(tokens.primary_color)

    def _write_column(self, frame, column_title: str | None, bullets: list[str], tokens: ThemeTokens) -> None:
        frame.clear()
        start_index = 0
        if column_title:
            title_p = frame.paragraphs[0]
            title_p.text = column_title
            title_p.font.name = tokens.heading_font
            title_p.font.size = Pt(16)
            title_p.font.bold = True
            title_p.font.color.rgb = rgb(tokens.accent_color)
            start_index = 1
        self._write_bullets(frame, bullets, tokens, start_index=start_index)

    def _write_bullets(self, frame, bullets: list[str], tokens: ThemeTokens, start_index: int = 0, font_color: str | None = None) -> None:
        if not bullets:
            return
        frame.word_wrap = True
        while len(frame.paragraphs) <= start_index:
            frame.add_paragraph()

        for index, bullet in enumerate(bullets):
            target_index = start_index + index
            paragraph = frame.paragraphs[target_index] if target_index < len(frame.paragraphs) else frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.name = tokens.body_font
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = rgb(font_color or tokens.body_color)
            paragraph.bullet = True

    def _style_shape_text(self, shape, font_name: str, font_size: int, font_color: str, *, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
        paragraph = shape.text_frame.paragraphs[0]
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = rgb(font_color)
        paragraph.font.bold = bold
        paragraph.alignment = align

    def _style_table_cell(self, cell, font_name: str, font_size: int, text_color: str, fill_color: str) -> None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(fill_color)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = rgb(text_color)

    def _fill_shape(self, shape, fill_color: str, line_color: str) -> None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill_color)
        shape.line.color.rgb = rgb(line_color)

    def _resolve_image_source(self, source: str) -> tuple[Path | None, str | None, Path | None]:
        if source.startswith(("http://", "https://")):
            try:
                response = requests.get(source, timeout=15)
                response.raise_for_status()
            except requests.RequestException as exc:
                return None, f"Failed to download image '{source}': {exc}", None
            suffix = os.path.splitext(source)[1] or ".img"
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
                temp_file.write(response.content)
                temp_path = Path(temp_file.name)
            return temp_path, None, temp_path

        path = Path(source)
        if path.exists():
            return path, None, None
        return None, f"Image path not found: {source}", None

    def _add_image(self, slide, image_path: Path) -> None:
        image = PptxImage.from_file(str(image_path))
        image_width_px, image_height_px = image.size
        aspect_ratio = image_width_px / image_height_px
        box_left = Inches(0.9)
        box_top = Inches(1.8)
        box_width = Inches(7.8)
        box_height = Inches(4.2)
        if (box_width / box_height) > aspect_ratio:
            height = box_height
            width = height * aspect_ratio
        else:
            width = box_width
            height = width / aspect_ratio
        left = box_left + (box_width - width) / 2
        top = box_top + (box_height - height) / 2
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    def _subtitle_text(self, slide_spec: SlideSpec) -> str | None:
        if isinstance(slide_spec, (TitleSlideSpec, SectionSlideSpec)):
            return slide_spec.subtitle
        return None


def required_template_roles(slide_spec: SlideSpec) -> set[str]:
    required = {"title"}
    if isinstance(slide_spec, (BulletsSlideSpec, TimelineSlideSpec, QuoteSlideSpec, ComparisonSlideSpec, SummarySlideSpec, TableSlideSpec, TwoColumnSlideSpec)):
        required.add("body")
    if isinstance(slide_spec, ImageSlideSpec):
        required.add("image")
    return required


def body_lines_for_slide(slide_spec: SlideSpec) -> list[str]:
    if isinstance(slide_spec, BulletsSlideSpec):
        return slide_spec.bullets
    if isinstance(slide_spec, TwoColumnSlideSpec):
        lines: list[str] = []
        if slide_spec.left_title:
            lines.append(slide_spec.left_title)
        lines.extend(f"- {item}" for item in slide_spec.left_bullets)
        if slide_spec.right_title:
            lines.append("")
            lines.append(slide_spec.right_title)
        lines.extend(f"- {item}" for item in slide_spec.right_bullets)
        return lines
    if isinstance(slide_spec, TimelineSlideSpec):
        return [f"{event.label}: {event.title}" + (f" - {event.detail}" if event.detail else "") for event in slide_spec.events]
    if isinstance(slide_spec, QuoteSlideSpec):
        details = [slide_spec.quote]
        byline = " / ".join(part for part in [slide_spec.attribution, slide_spec.source] if part)
        if byline:
            details.append(byline)
        return details
    if isinstance(slide_spec, ComparisonSlideSpec):
        lines = [slide_spec.left.title]
        lines.extend(f"- {item}" for item in slide_spec.left.bullets)
        lines.append("")
        lines.append(slide_spec.right.title)
        lines.extend(f"- {item}" for item in slide_spec.right.bullets)
        return lines
    if isinstance(slide_spec, SummarySlideSpec):
        lines = ["Key Points"]
        lines.extend(f"- {item}" for item in slide_spec.key_points)
        if slide_spec.next_steps:
            lines.append("")
            lines.append("Next Steps")
            lines.extend(f"- {item}" for item in slide_spec.next_steps)
        return lines
    if isinstance(slide_spec, TableSlideSpec):
        return [" | ".join(slide_spec.headers), *[" | ".join(row) for row in slide_spec.rows]]
    if isinstance(slide_spec, ImageSlideSpec):
        return slide_spec.bullets or ([slide_spec.caption] if slide_spec.caption else [])
    return []


def resolve_theme_tokens(theme: ThemeSpec) -> ThemeTokens:
    if theme.preset not in PRESET_THEMES:
        raise RenderError(f"Unsupported theme preset: {theme.preset}")
    preset = PRESET_THEMES[theme.preset]
    custom = theme.custom
    if custom is None:
        return preset
    return ThemeTokens(
        primary_color=custom.primary_color or preset.primary_color,
        accent_color=custom.accent_color or preset.accent_color,
        background_color=custom.background_color or preset.background_color,
        body_color=preset.body_color,
        muted_color=preset.muted_color,
        heading_font=custom.heading_font or preset.heading_font,
        body_font=custom.body_font or preset.body_font,
        cover_style=custom.cover_style or preset.cover_style,
    )


def rgb(color: str) -> RGBColor:
    return RGBColor.from_string(color.removeprefix("#"))


def lighten(color: str) -> str:
    rgb_values = tuple(int(color.removeprefix("#")[index : index + 2], 16) for index in range(0, 6, 2))
    mixed = tuple(min(255, int(channel + (255 - channel) * 0.78)) for channel in rgb_values)
    return "#" + "".join(f"{channel:02X}" for channel in mixed)
